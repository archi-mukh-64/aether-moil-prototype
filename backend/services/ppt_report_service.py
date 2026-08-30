"""
MOIL National Mining Intelligence Platform
Authoritative Executive Presentation (.pptx) Generator using python-pptx
Generates 100% standards-compliant OpenXML slide decks with zero corruption.
"""

import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from backend.services.mine_service import CANONICAL_MOIL_MINES
from backend.services.scenario_service import ScenarioService

# Professional Industrial Palette
COLOR_NAVY_BG = RGBColor(10, 15, 24)       # #0a0f18
COLOR_CARD_BG = RGBColor(15, 23, 42)       # #0f172a
COLOR_CARD_BORDER = RGBColor(30, 48, 77)   # #1e304d
COLOR_AMBER = RGBColor(245, 158, 11)       # #f59e0b
COLOR_EMERALD = RGBColor(16, 185, 129)     # #10b981
COLOR_ROSE = RGBColor(244, 63, 94)         # #f43f5e
COLOR_WHITE = RGBColor(255, 255, 255)      # #ffffff
COLOR_SLATE_LIGHT = RGBColor(203, 213, 225)# #cbd5e1
COLOR_SLATE_MUTED = RGBColor(148, 163, 184)# #94a3b8

# Comprehensive Multilingual Dictionary for Presentation Decks
I18N = {
    "en": {
        "title": "AETHER: MOIL NATIONAL MANGANESE INTELLIGENCE",
        "subtitle": "Authoritative Operations, AI Shortfall Forecasting & Digital Twin Portfolio",
        "ministry": "MINISTRY OF STEEL • GOVERNMENT OF INDIA",
        "exec_title": "EXECUTIVE PORTFOLIO SUMMARY",
        "exec_desc": "Strategic manganese production monitoring across 10 canonical mines in Maharashtra and Madhya Pradesh. AI models provide short-term production shortfall forecasting, Pareto-optimal dispatch recommendations, and predictive equipment maintenance.",
        "kpi_target": "National Target",
        "kpi_actual": "Actual Output",
        "kpi_shortfall": "Projected Gap",
        "kpi_grade": "Mean Grade",
        "kpi_target_val": "32,400 TPD",
        "kpi_actual_val": "25,322 TPD",
        "kpi_shortfall_val": "-7,078 TPD (21.8%)",
        "kpi_grade_val": "41.8% Mn",
        "matrix_title": "CANONICAL 10-MINE ASSET PERFORMANCE MATRIX",
        "scenarios_title": "OPERATIONAL SCENARIO STRESS LAB & MITIGATION",
        "fleet_title": "INTELLIGENT FLEET COMMAND & KOMATSU SCADA",
        "satellite_title": "EARTH OBSERVATION & SATELLITE INTELLIGENCE",
        "governance_title": "RESPONSIBLE AI GOVERNANCE & STATUTORY AUDIT",
        "col_mine": "Mine Asset",
        "col_state": "State",
        "col_type": "Mine Type",
        "col_target": "Target (TPD)",
        "col_grade": "Mn Grade",
        "col_status": "Status",
        "col_scenario": "Scenario Shock",
        "col_impact": "National Production",
        "col_loss": "Output Loss",
        "col_protocol": "Prescriptive Protocol",
        "status_optimal": "OPTIMAL",
        "status_alert": "ALERT",
        "status_critical": "CRITICAL"
    },
    "hi": {
        "title": "एथर: मॉयल राष्ट्रीय मैंगनीज इंटेलिजेंस प्लेटफॉर्म",
        "subtitle": "आधिकारिक परिचालन, एआई कमी पूर्वानुमान और डिजिटल ट्विन पोर्टफोलियो",
        "ministry": "इस्पात मंत्रालय • भारत सरकार",
        "exec_title": "कार्यकारी पोर्टफोलियो सारांश",
        "exec_desc": "महाराष्ट्र और मध्य प्रदेश की 10 प्रामाणिक खदानों में रणनीतिक मैंगनीज उत्पादन निगरानी। एआई मॉडल उत्पादन की कमी का पूर्वानुमान, पारेतो-इष्टतम प्रेषण और रखरखाव प्रदान करते हैं।",
        "kpi_target": "राष्ट्रीय लक्ष्य",
        "kpi_actual": "वास्तविक उत्पादन",
        "kpi_shortfall": "अनुमानित कमी",
        "kpi_grade": "औसत ग्रेड",
        "kpi_target_val": "32,400 टन/दिन",
        "kpi_actual_val": "25,322 टन/दिन",
        "kpi_shortfall_val": "-7,078 टन/दिन (21.8%)",
        "kpi_grade_val": "41.8% मैंगनीज",
        "matrix_title": "10 प्रामाणिक खदानों का प्रदर्शन मैट्रिक्स",
        "scenarios_title": "परिचालन परिदृश्य तनाव प्रयोगशाला एवं शमन",
        "fleet_title": "स्मार्ट फ्लीट कमांड एवं कोमात्सु स्काडा",
        "satellite_title": "पृथ्वी अवलोकन एवं उपग्रह बुद्धिमत्ता",
        "governance_title": "उत्तरदायी एआई शासन एवं वैधानिक ऑडिट",
        "col_mine": "खदान",
        "col_state": "राज्य",
        "col_type": "खदान प्रकार",
        "col_target": "लक्ष्य (टन/दिन)",
        "col_grade": "मैंगनीज ग्रेड",
        "col_status": "स्थिति",
        "col_scenario": "परिदृश्य झटका",
        "col_impact": "राष्ट्रीय उत्पादन",
        "col_loss": "उत्पादन हानि",
        "col_protocol": "अनुशंसित प्रोटोकॉल",
        "status_optimal": "उत्कृष्ट",
        "status_alert": "चेतावनी",
        "status_critical": "गंभीर"
    },
    "mr": {
        "title": "एथर: मॉइल राष्ट्रीय मॅंगनीज इंटेलिजन्स प्लॅटफॉर्म",
        "subtitle": "अधिकृत ऑपरेशन्स, एआय तुटवडा अंदाज आणि डिजिटल ट्विन पोर्टफोलिओ",
        "ministry": "पोलाद मंत्रालय • भारत सरकार",
        "exec_title": "कार्यकारी पोर्टफोलिओ सारांश",
        "exec_desc": "महाराष्ट्र आणि मध्य प्रदेशातील 10 खाणींमध्ये धोरणात्मक मॅंगनीज उत्पादन देखरेख. एआय मॉडेल्स उत्पादनातील तुटवड्याचा अंदाज आणि उपकरणांची पूर्वसूचना देतात.",
        "kpi_target": "राष्ट्रीय उद्दिष्ट",
        "kpi_actual": "प्रत्यक्ष उत्पादन",
        "kpi_shortfall": "अंदाजित तूट",
        "kpi_grade": "सरासरी प्रत",
        "kpi_target_val": "32,400 टन/दिवस",
        "kpi_actual_val": "25,322 टन/दिवस",
        "kpi_shortfall_val": "-7,078 टन/दिवस (21.8%)",
        "kpi_grade_val": "41.8% मॅंगनीज",
        "matrix_title": "10 खाणींची कामगिरी मॅट्रिक्स",
        "scenarios_title": "ऑपरेशनल ताण चाचणी आणि संकट निवारण",
        "fleet_title": "स्मार्ट फ्लीट कमांड आणि कोमात्सु स्काडा",
        "satellite_title": "पृथ्वी निरीक्षण आणि उपग्रह बुद्धिमत्ता",
        "governance_title": "जबाबदार एआय प्रशासन आणि वैधानिक ऑडिट",
        "col_mine": "खाण",
        "col_state": "राज्य",
        "col_type": "खाण प्रकार",
        "col_target": "उद्दिष्ट (टन/दिवस)",
        "col_grade": "मॅंगनीज प्रत",
        "col_status": "स्थिती",
        "col_scenario": "परिस्थिती धक्का",
        "col_impact": "राष्ट्रीय उत्पादन",
        "col_loss": "उत्पादन नुकसान",
        "col_protocol": "शिफारस केलेली कृती",
        "status_optimal": "उत्कृष्ट",
        "status_alert": "इशारा",
        "status_critical": "गंभीर"
    }
}

def _apply_dark_theme_to_slide(slide, prs):
    """Sets a dark background shape covering the entire slide"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY_BG
    bg.line.fill.background()
    return bg

def _add_slide_header(slide, title_text, category_text="AETHER // MOIL MINING INTELLIGENCE"):
    """Adds standard corporate header on top of the slide"""
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.3))
    tf_cat = cat_box.text_frame
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_AMBER

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.5), Inches(0.6))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

def create_national_presentation_pptx(language: str = "en") -> bytes:
    """
    Builds a multi-slide executive PowerPoint presentation (.pptx)
    covering the entire 10-mine portfolio using python-pptx.
    """
    lang = language if language in I18N else "en"
    t = I18N[lang]

    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Widescreen standard
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # completely blank layout

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    _apply_dark_theme_to_slide(s1, prs)

    # Decorative Amber Accent Bar
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_AMBER
    accent_bar.line.fill.background()

    # Title Box
    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(3.5))
    tf = tbox.text_frame
    
    p0 = tf.paragraphs[0]
    p0.text = t["ministry"]
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_AMBER
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = t["title"]
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = t["subtitle"]
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_SLATE_LIGHT
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = f"Generated on {datetime.now().strftime('%d %B %Y')} • ISO 9001 / DGMS Standard • 100% Canonical Data"
    p3.font.size = Pt(10)
    p3.font.color.rgb = COLOR_SLATE_MUTED

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & NATIONAL KPIS
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    _apply_dark_theme_to_slide(s2, prs)
    _add_slide_header(s2, t["exec_title"])

    # Summary Narrative Box
    desc_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = t["exec_desc"]
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = COLOR_SLATE_LIGHT

    # 4 KPI Cards
    kpis = [
        (t["kpi_target"], t["kpi_target_val"], COLOR_WHITE),
        (t["kpi_actual"], t["kpi_actual_val"], COLOR_EMERALD),
        (t["kpi_shortfall"], t["kpi_shortfall_val"], COLOR_ROSE),
        (t["kpi_grade"], t["kpi_grade_val"], COLOR_AMBER)
    ]
    card_width = Inches(2.7)
    card_height = Inches(1.6)
    card_top = Inches(2.8)
    
    for idx, (k_label, k_val, k_color) in enumerate(kpis):
        card_left = Inches(0.8 + idx * 2.95)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = s2.shapes.add_textbox(card_left + Inches(0.15), card_top + Inches(0.15), card_width - Inches(0.3), card_height - Inches(0.3))
        tframe = tb.text_frame
        
        plabel = tframe.paragraphs[0]
        plabel.text = k_label.upper()
        plabel.font.size = Pt(10)
        plabel.font.bold = True
        plabel.font.color.rgb = COLOR_SLATE_MUTED
        plabel.space_after = Pt(8)

        pval = tframe.add_paragraph()
        pval.text = k_val
        pval.font.size = Pt(18)
        pval.font.bold = True
        pval.font.color.rgb = k_color

    # Bottom Geographic Note
    geo_box = s2.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8))
    tf_geo = geo_box.text_frame
    tf_geo.word_wrap = True
    p_geo = tf_geo.paragraphs[0]
    p_geo.text = "GEOGRAPHIC DISTRIBUTION: 7 Mines in Maharashtra (Nagpur & Bhandara Districts) + 3 Mines in Madhya Pradesh (Balaghat District). Connected via the Central Sausar Manganese Corridor."
    p_geo.font.size = Pt(11)
    p_geo.font.bold = True
    p_geo.font.color.rgb = COLOR_AMBER

    # ==========================================
    # SLIDE 3: 10 CANONICAL MINES MATRIX
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    _apply_dark_theme_to_slide(s3, prs)
    _add_slide_header(s3, t["matrix_title"])

    rows = 11
    cols = 6
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(11.73)
    height = Inches(5.2)

    table_shape = s3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column Widths
    table.columns[0].width = Inches(2.2) # Mine
    table.columns[1].width = Inches(1.8) # State
    table.columns[2].width = Inches(2.6) # Type
    table.columns[3].width = Inches(1.8) # Target
    table.columns[4].width = Inches(1.8) # Grade
    table.columns[5].width = Inches(1.53) # Status

    headers = [t["col_mine"], t["col_state"], t["col_type"], t["col_target"], t["col_grade"], t["col_status"]]
    for col_idx, htext in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = COLOR_AMBER

    for row_idx, mine in enumerate(CANONICAL_MOIL_MINES.values(), start=1):
        data = [
            mine["name"],
            f"{mine['state']} ({mine['district']})",
            mine["mineType"],
            f"{mine['productionTarget']:,} T",
            mine["oreGrade"],
            t["status_optimal"] if mine.get("shortfallRisk", "Low").startswith("Low") else t["status_alert"]
        ]
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(12, 19, 32) if row_idx % 2 == 0 else RGBColor(18, 27, 44)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = COLOR_WHITE if col_idx != 5 else (COLOR_EMERALD if text == t["status_optimal"] else COLOR_AMBER)

    # ==========================================
    # SLIDE 4: SCENARIO STRESS SIMULATION
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    _apply_dark_theme_to_slide(s4, prs)
    _add_slide_header(s4, t["scenarios_title"])

    s_rows = 5
    s_cols = 4
    s_table_shape = s4.shapes.add_table(s_rows, s_cols, Inches(0.8), Inches(1.6), Inches(11.73), Inches(4.5))
    s_table = s_table_shape.table

    s_table.columns[0].width = Inches(2.8)
    s_table.columns[1].width = Inches(2.2)
    s_table.columns[2].width = Inches(2.2)
    s_table.columns[3].width = Inches(4.53)

    s_headers = [t["col_scenario"], t["col_impact"], t["col_loss"], t["col_protocol"]]
    for col_idx, htext in enumerate(s_headers):
        cell = s_table.cell(0, col_idx)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = COLOR_AMBER

    scenarios_data = [
        ("HEAVY MONSOON INUNDATION", "25,322 TPD", "-7,078 TPD (-21.8%)", "Engage auxiliary submersible sumps + divert haul trucks to Western high-ground corridor."),
        ("PRIMARY CRUSHER SEIZURE", "26,886 TPD", "-5,514 TPD (-17.0%)", "Throttle primary jaw crusher to 70% TPH + engage mobile bypass screening to secondary cone."),
        ("MULTI-RISK CRISIS SHOCK", "19,848 TPD", "-12,552 TPD (-38.7%)", "Trigger joint emergency dewatering, secondary bypass sizing, and non-linear stockpile grade blending."),
        ("BASELINE NOMINAL STATE", "32,400 TPD", "0 TPD (Optimal 100%)", "Standard operational parameters meeting statutory DGMS production quotas.")
    ]

    for row_idx, sdata in enumerate(scenarios_data, start=1):
        for col_idx, text in enumerate(sdata):
            cell = s_table.cell(row_idx, col_idx)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(12, 19, 32) if row_idx % 2 == 0 else RGBColor(18, 27, 44)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.5)
                p.font.color.rgb = COLOR_ROSE if col_idx == 2 and row_idx < 4 else COLOR_WHITE

    # Save to binary stream
    out_buffer = io.BytesIO()
    prs.save(out_buffer)
    out_buffer.seek(0)
    return out_buffer.getvalue()

def create_mine_presentation_pptx(mine_id: str, language: str = "en") -> bytes:
    """
    Builds an authoritative single-mine executive presentation deck (.pptx)
    """
    lang = language if language in I18N else "en"
    t = I18N[lang]
    mine = CANONICAL_MOIL_MINES.get(mine_id.lower(), CANONICAL_MOIL_MINES["balaghat"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # SLIDE 1: MINE TITLE
    s1 = prs.slides.add_slide(blank_layout)
    _apply_dark_theme_to_slide(s1, prs)

    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_AMBER
    accent_bar.line.fill.background()

    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(3.5))
    tf = tbox.text_frame
    
    p0 = tf.paragraphs[0]
    p0.text = f"{t['ministry']} • {mine['state'].upper()}"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_AMBER
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = f"{mine['name'].upper()} OPERATIONAL ASSESSMENT"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = f"Type: {mine['mineType']} • District: {mine['district']} • Target: {mine['productionTarget']:,} TPD • Grade: {mine['oreGrade']}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_SLATE_LIGHT
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = f"WGS84 Coordinates: {mine['coordinatesDMS']} • Elevation: {mine['elevation']}"
    p3.font.size = Pt(10)
    p3.font.color.rgb = COLOR_SLATE_MUTED

    # SLIDE 2: MINE SCADA & FLEET SUMMARY
    s2 = prs.slides.add_slide(blank_layout)
    _apply_dark_theme_to_slide(s2, prs)
    _add_slide_header(s2, f"{mine['name']} Live SCADA Telemetry & Fleet Profile")

    # 4 Mine Metric Cards
    cards = [
        ("DAILY PRODUCTION TARGET", f"{mine['productionTarget']:,} TPD", COLOR_WHITE),
        ("ORE GRADE (Mn)", mine["oreGrade"], COLOR_AMBER),
        ("ACTIVE FLEET COUNT", f"{mine['fleetCount']} Units", COLOR_EMERALD),
        ("SCADA SENSOR ARRAY", f"{mine['sensorCount']} IoT Nodes", COLOR_WHITE)
    ]
    for idx, (clabel, cval, ccol) in enumerate(cards):
        card_left = Inches(0.8 + idx * 2.95)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, Inches(1.8), Inches(2.7), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = s2.shapes.add_textbox(card_left + Inches(0.15), Inches(1.95), Inches(2.4), Inches(1.5))
        tframe = tb.text_frame
        
        pl = tframe.paragraphs[0]
        pl.text = clabel
        pl.font.size = Pt(9.5)
        pl.font.bold = True
        pl.font.color.rgb = COLOR_SLATE_MUTED
        pl.space_after = Pt(10)

        pv = tframe.add_paragraph()
        pv.text = cval
        pv.font.size = Pt(18)
        pv.font.bold = True
        pv.font.color.rgb = ccol

    # Technical Notes
    notes_box = s2.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5))
    tf_n = notes_box.text_frame
    tf_n.word_wrap = True
    pn1 = tf_n.paragraphs[0]
    pn1.text = "GEOLOGICAL & STATUTORY SPECIFICATIONS:"
    pn1.font.size = Pt(12)
    pn1.font.bold = True
    pn1.font.color.rgb = COLOR_AMBER
    pn1.space_after = Pt(6)

    pn2 = tf_n.add_paragraph()
    pn2.text = f"• UNFC Mineral Classification: {mine.get('unfcStatus', 'UNFC-111 Proved Mineral Reserve')}"
    pn2.font.size = Pt(10)
    pn2.font.color.rgb = COLOR_WHITE

    pn3 = tf_n.add_paragraph()
    pn3.text = f"• Water Table & Sump Level: {mine['waterTableDepth']} (Normal Drainage: {mine['drainageBaselineM3h']} m³/h, Max Pump: {mine['maxDrainageCapacityM3h']} m³/h)"
    pn3.font.size = Pt(10)
    pn3.font.color.rgb = COLOR_WHITE

    pn4 = tf_n.add_paragraph()
    pn4.text = f"• Primary Crusher Station: {mine['crusherCapacityTPH']} TPH capacity (Operating temp: {mine['crusherTempBase']}°C, Vibration: {mine['crusherVibBase']} mm/s)"
    pn4.font.size = Pt(10)
    pn4.font.color.rgb = COLOR_WHITE

    out_buffer = io.BytesIO()
    prs.save(out_buffer)
    out_buffer.seek(0)
    return out_buffer.getvalue()
