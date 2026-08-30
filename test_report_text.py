"""
Verification of multilingual PDF & PPTX text extraction and assertion test
"""

import io
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from backend.services.ppt_report_service import create_national_presentation_pptx
from backend.services.pdf_report_service import create_national_report_pdf

print("=" * 60)
print("TESTING PPTX MULTILINGUAL TEXT EXTRACTION")
print("=" * 60)

for lang, kw in [('en', 'AETHER'), ('hi', 'एथर'), ('mr', 'एथर')]:
    raw_pptx = create_national_presentation_pptx(lang)
    prs = Presentation(io.BytesIO(raw_pptx))
    extracted = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    extracted.append(paragraph.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        extracted.append(cell.text)
    txt = " ".join(extracted)
    has_kw = kw in txt
    print(f"PPTX [{lang}] -> Word count: {len(txt.split())} | Contains '{kw}': {has_kw}")
    assert has_kw, f"PPTX {lang} did not contain {kw}"

print("\n" + "=" * 60)
print("TESTING PDF MULTILINGUAL SIZE & UNIQUENESS")
print("=" * 60)

pdf_en = create_national_report_pdf('en')
pdf_hi = create_national_report_pdf('hi')
pdf_mr = create_national_report_pdf('mr')

print(f"PDF [EN] Size: {len(pdf_en)} bytes")
print(f"PDF [HI] Size: {len(pdf_hi)} bytes")
print(f"PDF [MR] Size: {len(pdf_mr)} bytes")

assert len(pdf_en) > 0 and len(pdf_hi) > 0 and len(pdf_mr) > 0
assert pdf_en != pdf_hi and pdf_en != pdf_mr and pdf_hi != pdf_mr

print("\nALL MULTILINGUAL PDF & PPTX ASSERTIONS PASSED!")
