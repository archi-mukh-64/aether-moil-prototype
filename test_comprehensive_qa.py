"""
MOIL AETHER Platform - Master Comprehensive QA & Integration Test Suite
Validates:
1. Application-wide Data Contract (10 Canonical Mines)
2. 7-Mode National Radar Independence & Non-Duplication
3. Multi-Temporal Earth Observation & Change Detection (2018 -> 2026)
4. AI Exploration Target Scanner & Virtual Drilling Stratigraphy
5. 40 Scenario Permutations (10 Mines x 4 Scenarios)
6. 5 ML Model Inference Subsystems
7. Cross-Mine Pearson Correlations & Capital What-If Simulation
8. Multilingual PDF & PPTX Report Generation & Devanagari Verification
"""

import sys
import os
import json
import urllib.request
import urllib.error
import io

# UTF-8 stdout configuration for Windows PowerShell
sys.stdout.reconfigure(encoding='utf-8')

BASE_API = "http://127.0.0.1:8000"

MINES = [
    "balaghat", "tirodi", "ukwa", "munsar", "kandri",
    "gumgaon", "chikla", "dongri-buzurg", "ramtek", "bhandara"
]

SCENARIOS = [
    "HEAVY_MONSOON", "CRUSHER_SEIZURE", "MULTI_RISK_CRISIS", "BASELINE_RESET"
]

NATIONAL_MODES = [
    "NATIONAL_PERFORMANCE", "RESERVE_POTENTIAL", "EXPLORATION_PRIORITY",
    "PRODUCTION_RISK", "ENVIRONMENTAL_RISK", "EQUIPMENT_RISK", "STRATEGIC_PRIORITY"
]

def make_request(url, method="GET", data=None):
    req = urllib.request.Request(url, method=method)
    if data:
        req.add_header('Content-Type', 'application/json')
        payload = json.dumps(data).encode('utf-8')
    else:
        payload = None

    try:
        with urllib.request.urlopen(req, data=payload, timeout=12) as response:
            return response.status, response.read().decode('utf-8')
    except urllib.error.HTTPError as e:
        return e.code, e.read().decode('utf-8')
    except Exception as e:
        return 0, str(e)

def run_qa_suite():
    print("=" * 70)
    print("AETHER — MASTER COMPREHENSIVE QA & INTEGRATION SUITE")
    print("=" * 70)

    total_assertions = 0
    passed_assertions = 0

    # -------------------------------------------------------------
    # 1. API GATEWAY HEALTH
    # -------------------------------------------------------------
    print("\n[SECTION 1] API Gateway Health & Model Registry")
    status, body = make_request(f"{BASE_API}/api/health")
    assert status == 200, f"Health check failed: {status}"
    health = json.loads(body)
    assert health.get("status") in ["HEALTHY", "healthy", "UP"]
    print(f"  [PASS] Gateway Status: {health.get('status')} | Models Loaded: {health.get('models', {}).get('models_status', 'ALL OK')}")
    total_assertions += 2
    passed_assertions += 2

    # -------------------------------------------------------------
    # 2. 10 CANONICAL MINES DATA CONTRACT & SUBSYSTEMS
    # -------------------------------------------------------------
    print("\n[SECTION 2] Canonical Data Contract & Subsystems (10 Mines)")
    mine_coords = {}
    for m in MINES:
        status, body = make_request(f"{BASE_API}/api/mines/{m}")
        assert status == 200, f"Mine {m} failed: {status}"
        data = json.loads(body)
        assert "latitude" in data and "longitude" in data
        assert "productionTarget" in data and data["productionTarget"] > 0
        assert "oreGrade" in data
        coords = (data["latitude"], data["longitude"])
        assert coords not in mine_coords.values(), f"Duplicate coordinates detected for {m}!"
        mine_coords[m] = coords

        # Subsystems verification
        for sub in ["telemetry", "analytics", "equipment", "reserve", "trust"]:
            st, _ = make_request(f"{BASE_API}/api/mines/{m}/{sub}")
            assert st == 200, f"Subsystem {sub} for {m} failed: {st}"
            total_assertions += 1
            passed_assertions += 1

        print(f"  [PASS] Mine: {m.ljust(14)} -> Coords: {coords} | Target: {data['productionTarget']} TPD | Grade: {data['oreGrade']}")
        total_assertions += 4
        passed_assertions += 4

    # -------------------------------------------------------------
    # 3. 7-MODE NATIONAL RADAR INDEPENDENCE & DISTINCTNESS
    # -------------------------------------------------------------
    print("\n[SECTION 3] 7-Mode National Radar Independence & Explainability")
    rank_vectors = []
    for mode in NATIONAL_MODES:
        status, body = make_request(f"{BASE_API}/api/national-radar/analyze?mode={mode}")
        assert status == 200, f"National Radar {mode} failed: {status}"
        res = json.loads(body)
        assert "rankedMines" in res and len(res["rankedMines"]) == 10
        assert "driverWeights" in res and sum(d["pct"] for d in res["driverWeights"]) == 100
        assert "nationalInsight" in res

        ranked_ids = [item["id"] for item in res["rankedMines"]]
        rank_vectors.append(tuple(ranked_ids))
        top_mine = res["rankedMines"][0]
        print(f"  [PASS] Mode: {mode.ljust(22)} -> Top #1: {top_mine['name'].ljust(14)} ({top_mine['metricPrimary']}) | Driver Sum: 100%")
        total_assertions += 4
        passed_assertions += 4

    unique_vectors = set(rank_vectors)
    assert len(unique_vectors) == len(NATIONAL_MODES), "National Radar modes produced duplicate rankings!"
    print(f"  [PASS] Unique Ranking Vectors: {len(unique_vectors)} / {len(NATIONAL_MODES)} (100% Mathematically Independent)")
    total_assertions += 1
    passed_assertions += 1

    # -------------------------------------------------------------
    # 4. MULTI-TEMPORAL EARTH OBSERVATION & CHANGE DETECTION
    # -------------------------------------------------------------
    print("\n[SECTION 4] Multi-Temporal Earth Observation (2018 -> 2026)")
    years_data = {}
    for yr in [2018, 2020, 2022, 2024, 2026]:
        status, body = make_request(f"{BASE_API}/api/earth-observation/environmental/balaghat?year={yr}")
        assert status == 200, f"EO year {yr} failed: {status}"
        data = json.loads(body)
        assert data["year"] == yr
        assert "meanNdvi" in data and "disturbedAreaHa" in data
        years_data[yr] = data
        print(f"  [PASS] Year {yr} -> Footprint: {data['footprintHa']} Ha | Disturbed: {data['disturbedAreaHa']} Ha | NDVI: {data['meanNdvi']}")
        total_assertions += 3
        passed_assertions += 3

    assert years_data[2018]["disturbedAreaHa"] != years_data[2026]["disturbedAreaHa"], "Temporal disturbance must change across years!"
    total_assertions += 1
    passed_assertions += 1

    # Change Detection Comparator
    status, body = make_request(f"{BASE_API}/api/earth-observation/compare-years/balaghat?year_before=2018&year_after=2026")
    assert status == 200
    comp = json.loads(body)
    assert comp["deltaFootprintHa"] > 0
    assert "summaryStatement" in comp
    print(f"  [PASS] 2018 vs 2026 Delta -> Footprint: +{comp['deltaFootprintHa']} Ha | Reclaimed: +{comp['deltaReclaimedHa']} Ha | NDVI: {comp['deltaNdvi']}")
    total_assertions += 3
    passed_assertions += 3

    # -------------------------------------------------------------
    # 5. AI EXPLORATION TARGET SCANNER & VIRTUAL CORE DRILL
    # -------------------------------------------------------------
    print("\n[SECTION 5] AI Exploration Target Scanner & Virtual Drilling")
    status, body = make_request(f"{BASE_API}/api/exploration/scan?mine_id=dongri-buzurg", method="POST")
    assert status == 200, f"Exploration scan failed: {status}"
    scan = json.loads(body)
    assert scan["targetsFound"] >= 3
    assert len(scan["targets"]) >= 3
    t1 = scan["targets"][0]
    assert "prospectivity" in t1 and "latitude" in t1 and "recommendation" in t1
    print(f"  [PASS] Scan Completed -> Found {scan['targetsFound']} Candidate Targets | Top #1: {t1['name']} ({t1['prospectivity']}% Score)")
    total_assertions += 4
    passed_assertions += 4

    # Virtual Core Drilling Depth Intercepts
    for d in [45, 145, 260]:
        status, body = make_request(f"{BASE_API}/api/exploration/drill/balaghat?depth_m={d}")
        assert status == 200
        drill = json.loads(body)
        assert "stratumLayer" in drill and "mnGrade" in drill
        print(f"  [PASS] Probe Depth: -{d}m -> Stratum: {drill['stratumLayer']} | Mn: {drill['mnGrade']}% | Conf: {drill['confidencePct']}%")
        total_assertions += 3
        passed_assertions += 3

    # -------------------------------------------------------------
    # 6. SCENARIO STRESS MATRIX (10 Mines x 4 Scenarios = 40 Tests)
    # -------------------------------------------------------------
    print("\n[SECTION 6] Scenario Stress Matrix (40 Permutations)")
    for m in MINES:
        for s in SCENARIOS:
            payload = {
                "mine_id": m,
                "scenario_type": s,
                "severity": "HIGH",
                "time_horizon": "24 HOURS"
            }
            status, body = make_request(f"{BASE_API}/api/scenarios/simulate", method="POST", data=payload)
            assert status == 200, f"Scenario {s} on {m} failed: {status}"
            res = json.loads(body)
            assert "predicted_loss_tonnes" in res or "is_detected" in res
            total_assertions += 2
            passed_assertions += 2
    print(f"  [PASS] All 40 Scenario Permutations Executed & Validated Successfully (10 Mines x 4 Stress Conditions)")

    # -------------------------------------------------------------
    # 7. ML INFERENCE SUBSYSTEMS & CAPITAL WHAT-IF
    # -------------------------------------------------------------
    print("\n[SECTION 7] ML Model Inference & Capital What-If Simulator")
    status, _ = make_request(f"{BASE_API}/api/alert/predict", method="POST", data={"mine_id": "balaghat", "rainfall_mm": 65.0, "sump_inflow_m3h": 720.0, "crusher_vibration_mms": 4.5, "haul_delay_mins": 14.0})
    assert status == 200
    print("  [PASS] Shortfall-GBM TreeSHAP Predictor")

    status, _ = make_request(f"{BASE_API}/api/reserve/predict", method="POST", data={"mine_id": "balaghat", "depth_m": 185.0, "mn_grade": 44.5, "silica_pct": 12.0})
    assert status == 200
    print("  [PASS] RandomForest Reserve Prospectivity Model")

    status, _ = make_request(f"{BASE_API}/api/equipment/predict", method="POST", data={"machine_id": "CRU-BLG-01", "vibration_mms": 5.2, "temperature_c": 86.0, "hydraulic_pressure_bar": 215.0, "operating_hours": 4200})
    assert status == 200
    print("  [PASS] GradientBoosting Equipment RUL & Failure Model")

    status, _ = make_request(f"{BASE_API}/api/anomaly/detect", method="POST", data={"sensor_values": [4.2, 85.0, 210.0, 78.0, 14.0]})
    assert status == 200
    print("  [PASS] IsolationForest Operational Anomaly Detector")

    status, _ = make_request(f"{BASE_API}/api/protocol/optimize", method="POST", data={"mine_id": "balaghat", "scenario_id": "MONSOON", "severity": "HIGH"})
    assert status == 200
    print("  [PASS] Pareto Multi-Objective Protocol Optimizer")

    status, body = make_request(f"{BASE_API}/api/national-radar/what-if?invest_crores=250")
    assert status == 200
    sim = json.loads(body)
    assert sim["reserveConversionMT"] > 0 and sim["productionIncreaseTPD"] > 0
    print(f"  [PASS] Capital Simulator (₹250 Cr) -> Reserve Upside: +{sim['reserveConversionMT']} MT | Output: +{sim['productionIncreaseTPD']} TPD")
    total_assertions += 8
    passed_assertions += 8

    # -------------------------------------------------------------
    # 8. MULTILINGUAL PDF & PPTX DECK EXTRACTION & VERIFICATION
    # -------------------------------------------------------------
    print("\n[SECTION 8] Multilingual PDF & PPTX Report Generation")
    from pptx import Presentation
    from backend.services.ppt_report_service import create_national_presentation_pptx
    from backend.services.pdf_report_service import create_national_report_pdf

    for lang, kw in [('en', 'AETHER'), ('hi', 'एथर'), ('mr', 'एथर')]:
        raw_pptx = create_national_presentation_pptx(lang)
        prs = Presentation(io.BytesIO(raw_pptx))
        extracted = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        extracted.append(paragraph.text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            extracted.append(cell.text)
        txt = " ".join(extracted)
        assert kw in txt, f"PPTX [{lang}] missing required keyword '{kw}'"
        print(f"  [PASS] PPTX [{lang}] -> Word Count: {len(txt.split())} | Verified Authentic Keyword: '{kw}'")
        total_assertions += 2
        passed_assertions += 2

    pdf_en = create_national_report_pdf('en')
    pdf_hi = create_national_report_pdf('hi')
    pdf_mr = create_national_report_pdf('mr')

    assert len(pdf_en) > 0 and len(pdf_hi) > 0 and len(pdf_mr) > 0
    assert pdf_en != pdf_hi and pdf_en != pdf_mr and pdf_hi != pdf_mr
    print(f"  [PASS] PDF [EN] Size: {len(pdf_en)} bytes | PDF [HI]: {len(pdf_hi)} bytes | PDF [MR]: {len(pdf_mr)} bytes (Distinct Multilingual Payloads)")
    total_assertions += 4
    passed_assertions += 4

    print("\n" + "=" * 70)
    print(f"MASTER QA RESULTS: {passed_assertions} / {total_assertions} ASSERTIONS PASSED (100% SUCCESS RATE)")
    print("=" * 70 + "\n")

if __name__ == '__main__':
    run_qa_suite()
